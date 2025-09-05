from langchain_community.document_loaders import (
    PDFPlumberLoader,
    UnstructuredPowerPointLoader,
    UnstructuredWordDocumentLoader,
    UnstructuredExcelLoader,
    CSVLoader,
    Docx2txtLoader,
    UnstructuredFileLoader
)
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
import os
import shutil
from typing import List
import hashlib
from pptx import Presentation  # Add this import for direct PPTX processing
import tempfile

# Initialize the vector database
FAISS_DB_PATH = "vectorstore/db_faiss"
faiss_db = None


# Step 1: Load Multiple Files of different types
def load_multiple_files(file_paths: List[str]):
    all_documents = []

    for file_path in file_paths:
        file_extension = os.path.splitext(file_path)[1].lower()

        try:
            if file_extension == '.pdf':
                loader = PDFPlumberLoader(file_path)
            elif file_extension in ['.pptx', '.ppt']:
                # Use a custom function to handle PPT/PPTX files
                documents = process_powerpoint_file(file_path)
                if documents:
                    all_documents.extend(documents)
                continue
            elif file_extension in ['.docx', '.doc']:
                # Try multiple loaders for DOCX files
                try:
                    loader = UnstructuredWordDocumentLoader(file_path)
                except:
                    loader = Docx2txtLoader(file_path)
            elif file_extension in ['.xlsx', '.xls']:
                loader = UnstructuredExcelLoader(file_path)
            elif file_extension == '.csv':
                loader = CSVLoader(file_path)
            else:
                print(f"Unsupported file type: {file_extension}")
                continue

            # Skip processing for PPT/PPTX files as we handled them above
            if file_extension in ['.pptx', '.ppt']:
                continue

            documents = loader.load()

            if not documents:
                print(f"No content found in {file_path}")
                continue

            # Add source information to each document
            for doc in documents:
                doc.metadata['source'] = file_path
                # Add a unique ID for each document
                doc.metadata['doc_id'] = hashlib.md5(f"{file_path}{doc.page_content[:100]}".encode()).hexdigest()

            all_documents.extend(documents)

        except Exception as e:
            print(f"Error loading {file_path}: {e}")
            continue

    return all_documents


# Custom function to process PowerPoint files
def process_powerpoint_file(file_path):
    """Process PowerPoint files and return documents"""
    try:
        # Try using UnstructuredPowerPointLoader first
        try:
            loader = UnstructuredPowerPointLoader(file_path, mode="elements")
            documents = loader.load()

            if documents and any(doc.page_content.strip() for doc in documents):
                for doc in documents:
                    doc.metadata['source'] = file_path
                    doc.metadata['doc_id'] = hashlib.md5(f"{file_path}{doc.page_content[:100]}".encode()).hexdigest()
                return documents
        except Exception as e:
            print(f"UnstructuredPowerPointLoader failed: {e}")

        # Fallback to direct python-pptx processing for PPT files
        return process_ppt_directly(file_path)

    except Exception as e:
        print(f"Error processing PowerPoint file {file_path}: {e}")
        return []


# Fallback function for direct PPT processing
def process_ppt_directly(file_path):
    """Process PPT files directly using python-pptx library"""
    try:
        # For .ppt files, we need to handle them differently than .pptx
        file_extension = os.path.splitext(file_path)[1].lower()

        if file_extension == '.ppt':
            # Convert .ppt to .pptx first (requires LibreOffice or similar)
            # If conversion isn't possible, we'll try to extract text using other methods
            return process_ppt_with_text_extraction(file_path)
        else:
            # Process .pptx files directly
            prs = Presentation(file_path)
            text_content = []

            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    # Check for tables
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if hasattr(cell, "text") and cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text.append(" | ".join(row_text))

                if slide_text:
                    text_content.append(f"Slide {i + 1}:\n" + "\n".join(slide_text))

            if text_content:
                # Create a single document with all content
                from langchain.schema import Document
                full_text = "\n\n".join(text_content)
                doc = Document(
                    page_content=full_text,
                    metadata={
                        'source': file_path,
                        'doc_id': hashlib.md5(f"{file_path}{full_text[:100]}".encode()).hexdigest()
                    }
                )
                return [doc]
            return []
    except Exception as e:
        print(f"Error in direct PPTX processing {file_path}: {e}")
        return []


# Alternative method for PPT files
def process_ppt_with_text_extraction(file_path):
    """Try to extract text from PPT files using alternative methods"""
    try:
        # Try using antiword or catdoc for text extraction
        # This is a fallback for older PPT formats
        import subprocess
        import tempfile

        # Create a temporary text file
        with tempfile.NamedTemporaryFile(suffix='.txt', delete=False) as temp_file:
            temp_path = temp_file.name

        try:
            # Try using catdoc (install with: sudo apt-get install catdoc)
            result = subprocess.run(['catdoc', file_path], capture_output=True, text=True)
            if result.returncode == 0 and result.stdout.strip():
                with open(temp_path, 'w', encoding='utf-8') as f:
                    f.write(result.stdout)

                # Load the extracted text
                loader = UnstructuredFileLoader(temp_path)
                documents = loader.load()

                if documents and any(doc.page_content.strip() for doc in documents):
                    for doc in documents:
                        doc.metadata['source'] = file_path
                        doc.metadata['doc_id'] = hashlib.md5(
                            f"{file_path}{doc.page_content[:100]}".encode()).hexdigest()
                    return documents
        except:
            pass

        # If catdoc failed, try using strings command
        try:
            result = subprocess.run(['strings', file_path], capture_output=True, text=True)
            if result.returncode == 0 and result.stdout.strip():
                # Filter out binary data and keep only text
                text_content = []
                for line in result.stdout.split('\n'):
                    if line.strip() and all(ord(c) < 128 for c in line):
                        text_content.append(line.strip())

                if text_content:
                    from langchain.schema import Document
                    full_text = "\n".join(text_content)
                    doc = Document(
                        page_content=full_text,
                        metadata={
                            'source': file_path,
                            'doc_id': hashlib.md5(f"{file_path}{full_text[:100]}".encode()).hexdigest()
                        }
                    )
                    return [doc]
        except:
            pass

        return []
    except Exception as e:
        print(f"Error in alternative PPT processing {file_path}: {e}")
        return []
    finally:
        # Clean up temporary file
        if os.path.exists(temp_path):
            os.unlink(temp_path)


# Step 2: Create Chunks with better chunking strategy
def create_chunks(documents):
    if not documents:
        return []

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        add_start_index=True,
        separators=["\n\n", "\n", ". ", " ", ""]  # Better splitting for various doc types
    )
    return text_splitter.split_documents(documents)


# Step 3: Setup Embeddings Model with better performance
def get_embedding_model():
    return HuggingFaceEmbeddings(
        model_name="sentence-transformers/all-mpnet-base-v2",
        model_kwargs={'device': 'cpu'},
        encode_kwargs={'normalize_embeddings': True}
    )


# Step 4: Create and save FAISS database from multiple files
def create_vector_db_from_multiple_files(file_paths: List[str], db_path: str):
    global faiss_db

    documents = load_multiple_files(file_paths)

    if not documents:
        raise Exception("No valid content found in the provided files.")

    text_chunks = create_chunks(documents)

    if not text_chunks:
        raise Exception("No text chunks could be created from the documents.")

    embeddings = get_embedding_model()
    faiss_db = FAISS.from_documents(text_chunks, embeddings)

    # Create directory if it doesn't exist
    os.makedirs(os.path.dirname(db_path), exist_ok=True)
    faiss_db.save_local(db_path)

    return faiss_db


# Update existing FAISS database with new files
def update_vector_db(new_file_paths: List[str], db_path: str):
    global faiss_db

    embeddings = get_embedding_model()

    if os.path.exists(db_path) and os.path.exists(os.path.join(db_path, "index.faiss")):
        try:
            # Load existing database
            faiss_db = FAISS.load_local(db_path, embeddings, allow_dangerous_deserialization=True)

            # Load new documents
            new_documents = load_multiple_files(new_file_paths)

            if new_documents:
                new_chunks = create_chunks(new_documents)
                if new_chunks:
                    faiss_db.add_documents(new_chunks)
        except Exception as e:
            print(f"Error loading existing database: {e}. Creating new one.")
            faiss_db = create_vector_db_from_multiple_files(new_file_paths, db_path)
    else:
        # Create new database if it doesn't exist
        os.makedirs(os.path.dirname(db_path), exist_ok=True)
        faiss_db = create_vector_db_from_multiple_files(new_file_paths, db_path)

    faiss_db.save_local(db_path)

    return faiss_db


# Load existing FAISS database
def load_vector_db(db_path):
    global faiss_db

    embeddings = get_embedding_model()

    if os.path.exists(db_path) and os.path.exists(os.path.join(db_path, "index.faiss")):
        try:
            db = FAISS.load_local(db_path, embeddings, allow_dangerous_deserialization=True)
            faiss_db = db
            return db
        except Exception as e:
            print(f"Error loading database: {e}")
            faiss_db = None
            return None
    else:
        print("Database path doesn't exist or is incomplete")
        faiss_db = None
        return None


# Clear the vector database
def clear_vector_db(db_path):
    global faiss_db

    if os.path.exists(db_path):
        try:
            shutil.rmtree(db_path)
            faiss_db = None
            return True
        except Exception as e:
            print(f"Error clearing database: {e}")
            return False
    return True


# Check if database is loaded
def is_db_loaded():
    global faiss_db

    if faiss_db is None:
        if os.path.exists(FAISS_DB_PATH) and os.path.exists(os.path.join(FAISS_DB_PATH, "index.faiss")):
            try:
                load_vector_db(FAISS_DB_PATH)
            except Exception as e:
                print(f"Error loading database: {e}")
                faiss_db = None
    return faiss_db is not None


# Get the database instance
def get_db():
    global faiss_db

    if faiss_db is None:
        if os.path.exists(FAISS_DB_PATH) and os.path.exists(os.path.join(FAISS_DB_PATH, "index.faiss")):
            try:
                faiss_db = load_vector_db(FAISS_DB_PATH)
            except Exception as e:
                print(f"Error loading database: {e}")
                faiss_db = None
    return faiss_db


# Enhanced retrieval with better similarity search
def enhanced_retrieval(query, k=4, threshold=0.6):
    if not is_db_loaded():
        raise Exception("Database not loaded. Please upload files first.")

    db = get_db()
    if db is None:
        raise Exception("Database not loaded. Please upload files first.")

    # First try similarity search
    results = db.similarity_search_with_score(query, k=k * 2)

    # Filter by score threshold
    filtered_results = [doc for doc, score in results if score >= threshold]

    # If we don't have enough results, try MMR for diversity
    if len(filtered_results) < k:
        try:
            mmr_results = db.max_marginal_relevance_search(query, k=k, fetch_k=k * 3)
            # Combine results
            combined = list({doc.metadata['doc_id']: doc for doc in filtered_results + mmr_results}.values())
            return combined[:k]
        except:
            # Fallback to regular similarity search
            return db.similarity_search(query, k=k)

    return filtered_results[:k]


# Get sources from retrieved documents
def get_sources(documents):
    sources = {}
    for doc in documents:
        if 'source' in doc.metadata:
            source_name = os.path.basename(doc.metadata['source'])
            if source_name not in sources:
                sources[source_name] = []
            sources[source_name].append(
                doc.page_content[:200] + "..." if len(doc.page_content) > 200 else doc.page_content)
    return sources